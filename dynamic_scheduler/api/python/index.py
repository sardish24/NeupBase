from fastapi import FastAPI, HTTPException
from pydantic import BaseModel
import pdfplumber
import docx
import pptx
import os
import aiofiles
import aiofiles.os
import tempfile
import httpx
import asyncio
from urllib.parse import urlparse

app = FastAPI()

class ExtractRequest(BaseModel):
    file_url: str
    document_type: str

def extract_text_from_file(temp_path: str, file_url: str, document_type: str) -> str:
    text_content = ""
    if document_type == "syllabus" or file_url.endswith(".pdf"):
        with pdfplumber.open(temp_path) as pdf:
            text_content = "\n".join([page.extract_text() for page in pdf.pages if page.extract_text()])
    elif file_url.endswith(".docx"):
        doc = docx.Document(temp_path)
        text_content = "\n".join([para.text for para in doc.paragraphs])
    elif file_url.endswith(".pptx"):
        presentation = pptx.Presentation(temp_path)
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text_content += shape.text + "\n"
    return text_content

@app.post(
    "/api/python/extract",
    responses={
        400: {"description": "Invalid file URL domain"},
        500: {"description": "Internal Server Error during extraction or LLM communication"}
    }
)
async def extract_topics(payload: ExtractRequest):
    try:
        temp_path = None
        # Step 1: Securely fetch file byte stream from Supabase bucket
        allowed_domains = ["supabase.co"]
        domain = urlparse(payload.file_url).netloc
        if not any(domain.endswith(d) for d in allowed_domains):
            raise HTTPException(status_code=400, detail="Invalid file URL domain")

        async with httpx.AsyncClient() as client:
            response = await client.get(payload.file_url)
            
            # Securely create a temporary file in the OS-designated temp folder asynchronously
            async with aiofiles.tempfile.NamedTemporaryFile(delete=False, prefix="doc_") as temp_file:
                await temp_file.write(response.content)
                temp_path = temp_file.name

        # Step 2: Protocol-dependent geometric and textual extraction
        text_content = await asyncio.to_thread(extract_text_from_file, temp_path, payload.file_url, payload.document_type)
        
        # Step 3: Large Language Model API Invocation (Google Gemini)
        gemini_api_key = os.getenv("GEMINI_API_KEY")
        gemini_prompt = f"""
        Analyze the following academic text. Extract the overarching subject topics, 
        their corresponding week numbers, and their academic format (lecture, tutorial, lab). 
        Return ONLY a strictly structured JSON array format. Do not include markdown formatting or conversational text.
        Text: {text_content[:15000]}  # Token limitation constraint
        """
        
        headers = {
            "Content-Type": "application/json"
        }
        data = {
            "contents": [{
                "parts": [{"text": gemini_prompt}]
            }],
            "generationConfig": {
                "responseMimeType": "application/json",
                "maxOutputTokens": 2048
            }
        }
        
        gemini_url = f"https://generativelanguage.googleapis.com/v1beta/models/gemini-3.5-flash:generateContent?key={gemini_api_key}"
        async with httpx.AsyncClient() as client:
            gemini_response = await client.post(gemini_url, headers=headers, json=data, timeout=60.0)
            response_json = gemini_response.json()
        
        try:
            import json
            extracted_text = response_json['candidates'][0]['content']['parts'][0]['text']
            parsed_topics = json.loads(extracted_text)
        except Exception:
            parsed_topics = response_json # fallback
            
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))
    finally:
        # Cleanup ephemeral file
        if 'temp_path' in locals() and temp_path and os.path.exists(temp_path):
            try:
                os.remove(temp_path)
            except Exception:
                pass

