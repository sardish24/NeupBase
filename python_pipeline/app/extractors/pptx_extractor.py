from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import requests
import tempfile
import os

class PPTXExtractor:
    @staticmethod
    def _extract_text_recursively(shapes, text_accumulator: list):
        """
        Recursively traverses grouped shapes to extract all embedded text runs.
        """
        for shape in shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                # Recurse into the nested shape hierarchy
                PPTXExtractor._extract_text_recursively(shape.shapes, text_accumulator)
            else:
                # Check for standard text frames
                if hasattr(shape, "text") and shape.text.strip():
                    text_accumulator.append(shape.text.strip())
                # Additionally handle tables explicitly embedded in slide geometry
                elif getattr(shape, "has_table", False):
                    for row in shape.table.rows:
                        row_data = [cell.text_frame.text.replace('\n', ' ').strip() for cell in row.cells]
                        text_accumulator.append(" | ".join(row_data))
        return text_accumulator

    @staticmethod
    def extract_text(file_url: str) -> str:
        extracted_content = []
        
        with tempfile.NamedTemporaryFile(delete=False, suffix=".pptx") as temp_pptx:
            response = requests.get(file_url, stream=True, timeout=45)
            response.raise_for_status()
            temp_pptx.write(response.content)
            temp_file_path = temp_pptx.name

        try:
            prs = Presentation(temp_file_path)
            for slide_num, slide in enumerate(prs.slides):
                extracted_content.append(f"\n=== SLIDE {slide_num + 1} ===")
                
                # Extract Slide Canvas Shapes via Deep Traversal
                slide_text = []
                PPTXExtractor._extract_text_recursively(slide.shapes, slide_text)
                extracted_content.extend(slide_text)
                
                # Extract Slide Notes (Highly relevant for academic lectures)
                if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                    notes = slide.notes_slide.notes_text_frame.text.strip()
                    if notes:
                        extracted_content.append(f"\n[Notes]:\n{notes}")
                        
        except Exception as e:
            raise RuntimeError(f"PPTX geometric extraction failed: {str(e)}")
        finally:
            if os.path.exists(temp_file_path):
                os.remove(temp_file_path)

        return "\n".join(extracted_content)
